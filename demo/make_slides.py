#!/usr/bin/env python3
"""Beta Hackathon submission deck — EXACTLY 3 slides (Team · Product · Demo).
Slide 3 embeds the demo video. Run with the venv that has python-pptx."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG = RGBColor(0x0A, 0x0C, 0x14); PANEL = RGBColor(0x10, 0x13, 0x1E)
FG = RGBColor(0xE8, 0xEA, 0xF2); MUTED = RGBColor(0x9A, 0xA0, 0xB6)
TEAL = RGBColor(0x2D, 0xD4, 0xA7); VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
BLUE = RGBColor(0x4A, 0xA8, 0xFF); ORANGE = RGBColor(0xFF, 0x78, 0x49)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
HERE = Path(__file__).resolve().parent


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.14), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False
    fr = s.shapes.add_textbox(Inches(0.6), Inches(7.02), Inches(12.3), Inches(0.4)).text_frame
    p = fr.paragraphs[0]; run = p.add_run()
    run.text = "Context Autopilot · Touchdown Labs · Built on EverMind + Butterbase (LIVE) · github.com/OCWC22/context-autopilot"
    run.font.size = Pt(10); run.font.color.rgb = MUTED
    return s


def tb(s, x, y, w, h):
    return s.shapes.add_textbox(x, y, w, h).text_frame


def txt(tf, text, size, color=FG, bold=False, align=PP_ALIGN.LEFT, space=8):
    p = tf.add_paragraph() if (tf.paragraphs[0].runs or tf.paragraphs[0].text) else tf.paragraphs[0]
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return p


def eyebrow(s, text, color=TEAL):
    txt(tb(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.5)), text, 15, color, bold=True)


def panel(s, x, y, w, h, accent):
    sp = s.shapes.add_shape(5, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = PANEL; sp.line.color.rgb = accent; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.28); tf.margin_top = Inches(0.22); tf.margin_right = Inches(0.22)
    return sp


# ── SLIDE 1 — TEAM INTRODUCTION ──
s = slide(); eyebrow(s, "SLIDE 1 — TEAM")
t = tb(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.4))
txt(t, "William Chen", 42, FG, bold=True)
txt(t, "Founder & CEO, Touchdown Labs — vendor-neutral inference-optimization research", 19, TEAL, bold=True)
tf = tb(s, Inches(0.6), Inches(2.9), Inches(12.1), Inches(2.0))
txt(tf, "Relevant experience", 16, VIOLET, bold=True, space=10)
for l in ["Built Edge AI and on-device / local AI before — shipping models that run close to the user.",
          "Touchdown Labs: research + tooling for cutting AI inference cost without losing quality."]:
    txt(tf, "•  " + l, 18, FG, space=10)
fp = panel(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.6), TEAL)
ff = fp.text_frame
txt(ff, "Team ⟷ problem fit", 16, TEAL, bold=True, space=8)
txt(ff, "Local & on-device AI is exactly my lane. This product makes local agents reuse structured context instead of re-paying the cloud to rediscover it — the on-device inference problem I already work on.", 16, FG)

# ── SLIDE 2 — PRODUCT OVERVIEW ──
s = slide(); eyebrow(s, "SLIDE 2 — PRODUCT")
txt(tb(s, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.8)), "Context Autopilot", 36, FG, bold=True)
txt(tb(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.9)),
    "One line: a $0 local engine that indexes ANY folder (code, docs, notes, projects) into reusable structured context — so Hermes, OpenClaw, or any local agent loads it instead of rediscovering everything.", 17, MUTED)
pp = panel(s, Inches(0.6), Inches(3.0), Inches(5.95), Inches(2.5), ORANGE)
txt(pp.text_frame, "Problem", 18, ORANGE, bold=True, space=8)
txt(pp.text_frame, "Local agents (Hermes, OpenClaw, Claude Code) rediscover your project every task — code, docs, notes — burning tokens reloading what they already saw (~90K to reload a 359 KB project).", 15, FG)
sp = panel(s, Inches(6.78), Inches(3.0), Inches(5.95), Inches(2.5), TEAL)
txt(sp.text_frame, "Solution", 18, TEAL, bold=True, space=8)
txt(sp.text_frame, "Index once into a reusable SKILL.md + DAG map; agents reuse it across tasks, frontier only for the hard step. EverMind = memory, Butterbase = state.", 15, FG)
txt(tb(s, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.0)),
    "Proof: same task 11,947 → 299 tokens (−98.7%, tests pass).  Real repo psf/requests: 147,794 → ~1,025 tokens (57×, −99.3%).", 17, BLUE, bold=True)

# ── SLIDE 3 — DEMO (embedded video) ──
s = slide(); eyebrow(s, "SLIDE 3 — DEMO")
txt(tb(s, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.7)),
    "Working product — watch the token meter drop (≤2 min)", 24, FG, bold=True)
vid = HERE / "autopilot-demo.mp4"; poster = HERE / "poster.png"
vw, vh = Inches(9.9), Inches(5.57)
vx = Emu(int((W - vw) / 2)); vy = Inches(1.7)
s.shapes.add_movie(str(vid), vx, vy, vw, vh,
                   poster_frame_image=str(poster) if poster.exists() else None, mime_type="video/mp4")

out = HERE / "autopilot-slides.pptx"
prs.save(str(out)); print("wrote", out, "slides:", len(prs.slides._sldIdLst))
